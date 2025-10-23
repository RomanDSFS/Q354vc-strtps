import os
os.environ["TRANSFORMERS_NO_TF"] = "1"
import uuid
import asyncio
import aiofiles
from pptx import Presentation
from functools import lru_cache
from transformers import pipeline
import fitz  # PyMuPDF
from config import UPLOAD_DIR

# ✅ Глобальное кеширование NLP-модели
@lru_cache(maxsize=1)
def get_nlp_model():
    return pipeline("zero-shot-classification", model="facebook/bart-large-mnli")

# Ключевые категории оценки
STARTUP_CATEGORIES = {
    "USP": ["unique selling proposition", "differentiation", "competitive advantage"],
    "Market": ["market size", "TAM", "SAM", "SOM", "growth potential"],
    "Business Model": ["revenue model", "monetization", "pricing strategy"],
    "Team": ["founders", "team experience", "leadership"],
    "Finance": ["funding", "financial projections", "investment"]
}

async def save_uploaded_file(file) -> str:
    """Асинхронное сохранение загруженного файла"""
    file_ext = file.filename.split('.')[-1]
    file_id = f"{uuid.uuid4()}.{file_ext}"
    file_path = os.path.join(UPLOAD_DIR, file_id)
    
    async with aiofiles.open(file_path, "wb") as buffer:
        while content := await file.read(8192):  # ✅ Увеличенный размер чанка
            await buffer.write(content)
    
    return file_path

async def extract_text(file_path: str, file_type: str) -> str:
    """Асинхронное извлечение текста из PDF и PPTX"""
    if file_type == "pdf":
        doc = fitz.open(file_path)
        return "\n".join([page.get_text() for page in doc])
    elif file_type == "pptx":
        prs = Presentation(file_path)
        return "\n".join([
            (slide.notes_text_frame.text if slide.notes_text_frame else "") + " " +
            " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            for slide in prs.slides
        ])
    else:
        return ""

async def analyze_text_async(text: str, keywords: list[str]):#(text, keywords):
    """Асинхронный анализ текста с кешированной моделью"""
    model = get_nlp_model()  # ✅ Используем кешированную модель
    return await asyncio.to_thread(model, text, candidate_labels=keywords)

async def analyze_startup_score_async(file_path: str, file_type: str) -> dict:
    """Асинхронный анализ питч-дека"""
    text = await extract_text(file_path, file_type)
    if not text:
        return {"error": "Unsupported or empty file"}

    tasks = [analyze_text_async(text, keywords) for _, keywords in STARTUP_CATEGORIES.items()]
    results = await asyncio.gather(*tasks)

    scores = {category: round(max(result["scores"]) * 20, 2) for category, result in zip(STARTUP_CATEGORIES.keys(), results)}
    total_score = sum(scores.values()) / 5 * 100

    return {"startup_score": round(total_score, 2), "details": scores}
